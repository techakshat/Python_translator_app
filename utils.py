import os
from docx import Document
from pptx import Presentation
from deep_translator import GoogleTranslator
from langdetect import detect, DetectorFactory
from datetime import datetime
import pandas as pd

from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import io
from pdf2image import convert_from_path
import pytesseract
from PIL import Image

# Ensure deterministic language detection
DetectorFactory.seed = 0

def detect_and_report_language(text):
    """Detects languages in a block of text and reports if multiple are found."""
    if not text.strip():
        return ""
    try:
        chunks = text.split('\n')
        detected_languages = set()
        for chunk in chunks:
            if chunk.strip():
                detected_languages.add(detect(chunk))
        if len(detected_languages) > 1:
            return "multi"
        elif len(detected_languages) == 1:
            return detected_languages.pop()
        else:
            return ""
    except Exception as e:
        print(f"Error detecting language: {e}") # Log the error
        return ""

def translate_text_chunk(text, dest_lang):
    """Translates a chunk of text."""
    if not text.strip():
        return text
    try:
        return GoogleTranslator(source='auto', target=dest_lang).translate(text)
    except Exception as e:
        print(f"Error translating text chunk: {e}") # Log the error
        return text

def translate_docx(file_path, dest_lang):
    """Translates a .docx file and returns the path to the new file."""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                para.text = translate_text_chunk(para.text, dest_lang)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    cell.text = translate_text_chunk(cell.text, dest_lang)
        new_file_path = os.path.splitext(file_path)[0] + f"_translated_{dest_lang}.docx"
        doc.save(new_file_path)
        return new_file_path
    except Exception as e:
        print(f"Error translating DOCX file {file_path}: {e}")
        raise

def translate_pptx(file_path, dest_lang):
    """Translates a .pptx file and returns the path to the new file."""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = translate_text_chunk(run.text, dest_lang)
        new_file_path = os.path.splitext(file_path)[0] + f"_translated_{dest_lang}.pptx"
        prs.save(new_file_path)
        return new_file_path
    except Exception as e:
        print(f"Error translating PPTX file {file_path}: {e}")
        raise

def translate_pdf(input_path, target_lang_code):
    """Translate text-based PDF. Returns (output_path, text_found)."""
    try:
        reader = PdfReader(input_path)
        writer = PdfWriter()
        text_found = False
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                text_found = True
                translated_text = translate_text_chunk(text, target_lang_code)
                packet = io.BytesIO()
                can = canvas.Canvas(packet, pagesize=letter)
                can.setFont("Helvetica", 12)
                y = 750
                for line in translated_text.split('\n'):
                    can.drawString(40, y, line[:1000])
                    y -= 15
                    if y < 40:
                        break
                can.save()
                packet.seek(0)
                translated_page = PdfReader(packet).pages[0]
                writer.add_page(translated_page)
            else:
                writer.add_page(page)
        output_path = input_path.replace(".pdf", f"_translated_{target_lang_code}.pdf")
        with open(output_path, "wb") as f:
            writer.write(f)
        return output_path, text_found
    except Exception as e:
        print(f"Error translating PDF file {input_path}: {e}")
        raise

def translate_pdf_ocr(input_path, target_lang_code):
    """OCR-based PDF translation for scanned/image PDFs."""
    try:
        images = convert_from_path(input_path)
        writer = PdfWriter()
        for img in images:
            text = pytesseract.image_to_string(img)
            translated_text = translate_text_chunk(text, target_lang_code)
            packet = io.BytesIO()
            can = canvas.Canvas(packet, pagesize=letter)
            can.setFont("Helvetica", 12)
            y = 750
            for line in translated_text.split('\n'):
                can.drawString(40, y, line[:1000])
                y -= 15
                if y < 40:
                    break
            can.save()
            packet.seek(0)
            translated_page = PdfReader(packet).pages[0]
            writer.add_page(translated_page)
        output_path = input_path.replace(".pdf", f"_translated_ocr_{target_lang_code}.pdf")
        with open(output_path, "wb") as f:
            writer.write(f)
        return output_path
    except Exception as e:
        print(f"Error translating PDF (OCR) file {input_path}: {e}")
        raise

def log_activity(activity_type, file_name, source_language, target_language):
    """Logs user activity to a CSV file."""
    log_file = "user_log.csv"
    log_entry = {
        "timestamp": [datetime.now().strftime("%Y-%m-%d %H:%M:%S")],
        "activity_type": [activity_type],
        "file_name": [file_name],
        "source_language": [source_language],
        "target_language": [target_language],
    }
    df = pd.DataFrame(log_entry)
    if os.path.exists(log_file):
        df.to_csv(log_file, mode="a", header=False, index=False)
    else:
        df.to_csv(log_file, index=False)